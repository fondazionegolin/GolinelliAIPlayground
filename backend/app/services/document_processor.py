"""
Multi-step document processor.

Pipeline:
  1. Extract text (PyMuPDF → PyPDF2 fallback, python-docx, plain text)
  2. Detect visual-heavy pages (charts, diagrams, images)
  3. Analyze visual pages with vision LLM (GPT-4o)
  4. Generate structured summary (summary, key_concepts, entities)
  5. Return DocumentAnalysis with ready-to-inject context string
"""

import io
import re
import json
import base64
import logging
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)

MAX_SUMMARY_SOURCE_CHARS = 12000
MAX_STRUCTURED_PREVIEW_CHARS = 4000
MAX_VISUAL_PAGES = 8
VISION_MODEL = "gpt-4o"
VISION_PROVIDER = "openai"


@dataclass
class DocumentSegment:
    text: str
    page: Optional[int] = None
    kind: str = "text"
    meta: dict = field(default_factory=dict)


@dataclass
class DocumentAnalysis:
    filename: str
    mime_type: str
    raw_text: str
    page_count: int
    summary: str
    key_concepts: list
    entities: list
    visual_descriptions: list
    rag_segments: list[DocumentSegment]
    structured_extract: str       # Ready-to-inject context block
    has_visual_content: bool
    processing_steps: list


class DocumentProcessor:
    def _stringify_cell(self, value) -> str:
        text = str(value).strip()
        if not text or text.lower() == "nan":
            return ""
        if len(text) > 200:
            return text[:197] + "..."
        return text

    def _segments_from_blocks(
        self,
        text: str,
        *,
        page: Optional[int] = None,
        kind: str = "text",
        target_chars: int = 1500,
    ) -> list[DocumentSegment]:
        blocks = [block.strip() for block in re.split(r"\n\s*\n", text or "") if block.strip()]
        if not blocks:
            return []

        segments: list[DocumentSegment] = []
        current_blocks: list[str] = []
        current_len = 0

        for block in blocks:
            block_len = len(block)
            if current_blocks and current_len + block_len + 2 > target_chars:
                segments.append(
                    DocumentSegment(
                        text="\n\n".join(current_blocks),
                        page=page,
                        kind=kind,
                    )
                )
                current_blocks = []
                current_len = 0
            current_blocks.append(block)
            current_len += block_len + 2

        if current_blocks:
            segments.append(
                DocumentSegment(
                    text="\n\n".join(current_blocks),
                    page=page,
                    kind=kind,
                )
            )

        return segments

    async def process(
        self,
        file_bytes: bytes,
        filename: str,
        mime_type: str,
        llm_service=None,
        analyze_visuals: bool = True,
    ) -> DocumentAnalysis:
        steps: list[str] = []

        # Step 1: Extract text
        raw_text, page_count, visual_page_indices, rag_segments = await self._extract_content(
            file_bytes, filename, mime_type, steps
        )

        # Step 2: Analyze visual pages
        visual_segments: list[DocumentSegment] = []
        has_visual = len(visual_page_indices) > 0

        if analyze_visuals and has_visual and llm_service:
            fn_lower = filename.lower()
            if mime_type == "application/pdf" or fn_lower.endswith(".pdf"):
                visual_segments = await self._analyze_pdf_visuals(
                    file_bytes, visual_page_indices, llm_service, steps
                )
            elif visual_page_indices == [-1]:
                visual_segments = await self._analyze_image_file(
                    file_bytes, filename, mime_type, llm_service, steps
                )
        visual_descriptions = [segment.text for segment in visual_segments]
        rag_segments.extend(visual_segments)

        # Step 3: Structured summary
        summary, key_concepts, entities = await self._generate_summary(
            raw_text, visual_descriptions, llm_service, steps
        )

        # Step 4: Build context string
        structured_extract = self._build_structured_extract(
            filename, summary, key_concepts, entities, visual_descriptions, raw_text
        )

        return DocumentAnalysis(
            filename=filename,
            mime_type=mime_type,
            raw_text=raw_text,
            page_count=page_count,
            summary=summary,
            key_concepts=key_concepts,
            entities=entities,
            visual_descriptions=visual_descriptions,
            rag_segments=rag_segments,
            structured_extract=structured_extract,
            has_visual_content=has_visual,
            processing_steps=steps,
        )

    # -------------------------------------------------------------------------
    # Text extraction
    # -------------------------------------------------------------------------

    async def _extract_content(
        self, file_bytes: bytes, filename: str, mime_type: str, steps: list
    ) -> tuple:
        """Returns (raw_text, page_count, visual_page_indices, rag_segments)."""
        fn_lower = filename.lower()

        if mime_type == "application/pdf" or fn_lower.endswith(".pdf"):
            return await self._extract_pdf(file_bytes, steps)

        if (
            mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or fn_lower.endswith(".docx")
        ):
            text, segments = self._extract_docx(file_bytes, steps)
            return text, 1, [], segments

        if (
            mime_type in (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint",
            )
            or fn_lower.endswith((".pptx", ".ppt"))
        ):
            text, segments = self._extract_pptx(file_bytes, steps)
            return text, 1, [], segments

        if (
            mime_type in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            )
            or fn_lower.endswith((".xlsx", ".xls"))
        ):
            text, segments = self._extract_xlsx(file_bytes, steps)
            return text, 1, [], segments

        if fn_lower.endswith(".csv") or mime_type in ("text/csv", "application/csv"):
            text, segments = self._extract_csv(file_bytes, steps)
            return text, 1, [], segments

        if mime_type.startswith("text/") or fn_lower.endswith((".txt", ".md")):
            try:
                text = file_bytes.decode("utf-8", errors="ignore")
            except Exception:
                text = ""
            steps.append(f"Testo piano estratto ({len(text)} caratteri)")
            return text, 1, [], self._segments_from_blocks(text)

        if mime_type.startswith("image/"):
            steps.append("File immagine: analisi visiva")
            return "", 1, [-1], []

        steps.append(f"Tipo file non supportato: {mime_type}")
        return "", 0, [], []

    def _extract_csv(self, file_bytes: bytes, steps: list) -> tuple[str, list[DocumentSegment]]:
        """Parse CSV and convert to readable markdown table + stats."""
        try:
            import pandas as pd
            import io as _io
            # Try common encodings
            for enc in ("utf-8", "latin-1", "cp1252"):
                try:
                    df = pd.read_csv(_io.BytesIO(file_bytes), encoding=enc)
                    break
                except Exception:
                    continue
            else:
                text = file_bytes.decode("utf-8", errors="ignore")
                steps.append("CSV letto come testo grezzo (encoding fallback)")
                return text, self._segments_from_blocks(text)

            text, segments = self._dataframe_to_text(df, steps, "CSV")
            return text, segments
        except Exception as e:
            steps.append(f"Errore lettura CSV ({e}), fallback testo grezzo")
            try:
                text = file_bytes.decode("utf-8", errors="ignore")
                return text, self._segments_from_blocks(text)
            except Exception:
                return "", []

    def _extract_xlsx(self, file_bytes: bytes, steps: list) -> tuple[str, list[DocumentSegment]]:
        """Parse XLSX and convert all sheets to readable text."""
        try:
            import pandas as pd
            import io as _io
            xl = pd.ExcelFile(_io.BytesIO(file_bytes))
            sheet_parts: list[str] = []
            segments: list[DocumentSegment] = []
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                sheet_text, sheet_segments = self._dataframe_to_text(df, steps, sheet_name)
                sheet_parts.append(sheet_text)
                for segment in sheet_segments:
                    segment.meta.setdefault("sheet_name", sheet_name)
                segments.extend(sheet_segments)
            text = "\n\n".join(sheet_parts)
            steps.append(f"XLSX estratto: {len(xl.sheet_names)} fogli")
            return text, segments
        except Exception as e:
            steps.append(f"Errore lettura XLSX: {e}")
            return "", []

    def _dataframe_to_text(self, df, steps: list, label: str) -> tuple[str, list[DocumentSegment]]:
        """Convert a DataFrame to a full readable text representation for RAG."""
        try:
            import pandas as pd

            rows, cols = df.shape
            col_names = list(df.columns)

            summary_lines: list[str] = []
            summary_lines.append(f"Tabella {label}: {rows} righe × {cols} colonne")
            summary_lines.append(f"Colonne: {', '.join(str(c) for c in col_names)}")

            # Numeric statistics for numeric columns
            num_cols = df.select_dtypes(include="number").columns.tolist()
            if num_cols:
                summary_lines.append("\n### Statistiche colonne numeriche")
                for col in num_cols[:10]:
                    s = df[col].dropna()
                    if len(s):
                        summary_lines.append(
                            f"- {col}: min={s.min():.4g}, max={s.max():.4g}, "
                            f"media={s.mean():.4g}, mediana={s.median():.4g}, "
                            f"tot={s.sum():.4g} ({len(s)} valori)"
                        )

            # Value counts for categorical columns
            cat_cols = df.select_dtypes(exclude="number").columns.tolist()
            if cat_cols:
                summary_lines.append("\n### Valori categorici principali")
                for col in cat_cols[:5]:
                    vc = df[col].value_counts().head(10)
                    summary_lines.append(f"- {col}: " + ", ".join(f"{v} ({c})" for v, c in vc.items()))

            normalized = df.where(pd.notnull(df), "")
            row_segments: list[DocumentSegment] = []
            batch_size = 50
            for start in range(0, rows, batch_size):
                end = min(start + batch_size, rows)
                batch_lines = [f"### {label} righe {start + 1}-{end}"]
                for row_idx in range(start, end):
                    row = normalized.iloc[row_idx]
                    cells = []
                    for col_name, value in row.items():
                        rendered = self._stringify_cell(value)
                        if rendered:
                            cells.append(f"{col_name}={rendered}")
                    if cells:
                        batch_lines.append(f"Riga {row_idx + 1}: " + " | ".join(cells))
                if len(batch_lines) > 1:
                    row_segments.append(
                        DocumentSegment(
                            text="\n".join(batch_lines),
                            kind="table_rows",
                            meta={"source_label": label, "row_start": start + 1, "row_end": end},
                        )
                    )

            segments = [
                DocumentSegment(
                    text="\n".join(summary_lines),
                    kind="table_summary",
                    meta={"source_label": label},
                )
            ] + row_segments

            text = "\n\n".join(segment.text for segment in segments if segment.text.strip())
            steps.append(f"{label} convertito: {rows} righe, {cols} colonne")
            return text, segments
        except Exception as e:
            steps.append(f"Errore conversione dataframe ({e})")
            fallback = str(df)
            return fallback, self._segments_from_blocks(fallback)

    async def _extract_pdf(self, file_bytes: bytes, steps: list) -> tuple:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            page_count = len(doc)
            parts: list[str] = []
            visual_pages: list[int] = []
            segments: list[DocumentSegment] = []

            for i in range(page_count):
                page = doc[i]
                page_text = (page.get_text() or "").strip()
                text_len = len(page_text.strip())

                n_images = len(page.get_images())
                n_drawings = len(page.get_drawings())
                table_count = 0
                if hasattr(page, "find_tables"):
                    try:
                        tables = page.find_tables()
                        table_count = len(getattr(tables, "tables", []) or [])
                    except Exception:
                        table_count = 0
                is_visual = bool(n_images or table_count or (n_drawings > 2 and text_len < 800))

                if is_visual and len(visual_pages) < MAX_VISUAL_PAGES:
                    visual_pages.append(i)

                page_header = f"[Pagina {i + 1}]"
                if page_text:
                    page_content = f"{page_header}\n{page_text}"
                    parts.append(page_content)
                    segments.append(
                        DocumentSegment(
                            text=page_content,
                            page=i + 1,
                            kind="page_text",
                            meta={"page_number": i + 1},
                        )
                    )

            doc.close()
            full_text = "\n\n".join(parts)
            steps.append(
                f"PDF estratto con PyMuPDF: {page_count} pagine, {len(segments)} pagine testuali, {len(visual_pages)} pagine visive"
            )
            return full_text, page_count, visual_pages, segments

        except ImportError:
            steps.append("PyMuPDF non disponibile, uso PyPDF2")
            return await self._extract_pdf_fallback(file_bytes, steps)
        except Exception as e:
            steps.append(f"Errore PyMuPDF ({e}), uso PyPDF2")
            return await self._extract_pdf_fallback(file_bytes, steps)

    async def _extract_pdf_fallback(self, file_bytes: bytes, steps: list) -> tuple:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            parts: list[str] = []
            segments: list[DocumentSegment] = []
            for idx, page in enumerate(reader.pages, 1):
                page_text = (page.extract_text() or "").strip()
                if not page_text:
                    continue
                page_content = f"[Pagina {idx}]\n{page_text}"
                parts.append(page_content)
                segments.append(
                    DocumentSegment(
                        text=page_content,
                        page=idx,
                        kind="page_text",
                        meta={"page_number": idx},
                    )
                )
            text = "\n\n".join(parts)
            steps.append(f"PDF estratto con PyPDF2: {len(reader.pages)} pagine")
            return text, len(reader.pages), [], segments
        except Exception as e:
            steps.append(f"Errore PyPDF2: {e}")
            return "", 0, [], []

    def _extract_docx(self, file_bytes: bytes, steps: list) -> tuple[str, list[DocumentSegment]]:
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            steps.append(f"DOCX estratto: {len(paragraphs)} paragrafi")
            return text, self._segments_from_blocks(text)
        except Exception as e:
            steps.append(f"Errore DOCX: {e}")
            return "", []

    def _extract_pptx(self, file_bytes: bytes, steps: list) -> tuple[str, list[DocumentSegment]]:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slide_texts: list[str] = []
            segments: list[DocumentSegment] = []
            for i, slide in enumerate(prs.slides, 1):
                parts: list[str] = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs if run.text.strip())
                        if line.strip():
                            parts.append(line.strip())
                if parts:
                    slide_text = f"## Slide {i}\n" + "\n".join(parts)
                    slide_texts.append(slide_text)
                    segments.append(
                        DocumentSegment(
                            text=slide_text,
                            page=i,
                            kind="slide_text",
                            meta={"slide_number": i},
                        )
                    )
            text = "\n\n".join(slide_texts)
            steps.append(f"PPTX estratto: {len(prs.slides)} slide, {len(text)} caratteri")
            return text, segments
        except Exception as e:
            steps.append(f"Errore PPTX: {e}")
            return "", []

    # -------------------------------------------------------------------------
    # Visual analysis
    # -------------------------------------------------------------------------

    async def _analyze_pdf_visuals(
        self, file_bytes: bytes, page_indices: list, llm_service, steps: list
    ) -> list[DocumentSegment]:
        descriptions: list[DocumentSegment] = []
        try:
            import fitz
            doc = fitz.open(stream=file_bytes, filetype="pdf")

            for page_idx in page_indices[:MAX_VISUAL_PAGES]:
                try:
                    page = doc[page_idx]
                    mat = fitz.Matrix(1.5, 1.5)   # ~108 DPI – good quality/size tradeoff
                    pix = page.get_pixmap(matrix=mat)
                    img_b64 = base64.b64encode(pix.tobytes("png")).decode("utf-8")

                    vision_resp = await llm_service.generate(
                        messages=[{
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": (
                                        "Descrivi il contenuto visivo di questa pagina: grafici, tabelle, "
                                        "diagrammi, immagini. Estrai tutti i dati numerici e le informazioni "
                                        "chiave. Sii preciso e conciso. Rispondi in italiano."
                                    ),
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {"url": f"data:image/png;base64,{img_b64}"},
                                },
                            ],
                        }],
                        provider=VISION_PROVIDER,
                        model=VISION_MODEL,
                        temperature=0.2,
                        max_tokens=500,
                    )
                    descriptions.append(
                        DocumentSegment(
                            text=f"[Pagina {page_idx + 1}] Analisi visiva: {vision_resp.content}",
                            page=page_idx + 1,
                            kind="visual_analysis",
                            meta={"page_number": page_idx + 1, "visual_analysis": True},
                        )
                    )
                    steps.append(f"Analizzata pagina visiva {page_idx + 1}")
                except Exception as e:
                    steps.append(f"Errore analisi visiva pagina {page_idx + 1}: {e}")

            doc.close()
        except ImportError:
            steps.append("PyMuPDF non disponibile per analisi visiva")
        except Exception as e:
            steps.append(f"Errore analisi visiva PDF: {e}")

        return descriptions

    async def _analyze_image_file(
        self, file_bytes: bytes, filename: str, mime_type: str, llm_service, steps: list
    ) -> list[DocumentSegment]:
        try:
            img_b64 = base64.b64encode(file_bytes).decode("utf-8")
            vision_resp = await llm_service.generate(
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Descrivi dettagliatamente questa immagine. "
                                "Se contiene grafici, tabelle o dati, estraili tutti. "
                                "Rispondi in italiano."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:{mime_type};base64,{img_b64}"},
                        },
                    ],
                }],
                provider=VISION_PROVIDER,
                model=VISION_MODEL,
                temperature=0.2,
                max_tokens=600,
            )
            steps.append("Immagine analizzata")
            return [
                DocumentSegment(
                    text=f"[{filename}] Analisi immagine: {vision_resp.content}",
                    kind="visual_analysis",
                    meta={"filename": filename, "visual_analysis": True},
                )
            ]
        except Exception as e:
            steps.append(f"Errore analisi immagine: {e}")
            return []

    # -------------------------------------------------------------------------
    # Summary generation
    # -------------------------------------------------------------------------

    async def _generate_summary(
        self, raw_text: str, visual_descriptions: list, llm_service, steps: list
    ) -> tuple:
        """Returns (summary, key_concepts, entities)."""
        if not raw_text and not visual_descriptions:
            return "Documento vuoto o non analizzabile.", [], []

        if not llm_service:
            words = raw_text.split()
            summary = " ".join(words[:80]) + "..." if len(words) > 80 else raw_text
            return summary, [], []

        content = raw_text[:MAX_SUMMARY_SOURCE_CHARS]
        if visual_descriptions:
            content += "\n\nCONTENUTO VISIVO:\n" + "\n".join(visual_descriptions[:3])

        try:
            response = await llm_service.generate(
                messages=[{
                    "role": "user",
                    "content": (
                        'Analizza questo documento e rispondi in JSON:\n'
                        '{"summary":"Riassunto sintetico (3-5 frasi)",'
                        '"key_concepts":["concetto1","concetto2"],'
                        '"entities":["entità1","entità2"]}\n\n'
                        f'DOCUMENTO:\n{content}\n\nRispondi SOLO con JSON valido.'
                    ),
                }],
                system_prompt=(
                    "Sei un esperto analista di documenti. "
                    "Estrai informazioni strutturate in JSON. "
                    "Rispondi SOLO con JSON valido, nessun altro testo."
                ),
                temperature=0.2,
                max_tokens=600,
            )

            text = response.content.strip()
            # Strip markdown code fences
            text = re.sub(r"^```\w*\n?", "", text)
            text = re.sub(r"\n?```$", "", text)

            data = json.loads(text)
            summary = str(data.get("summary", ""))
            key_concepts = [str(c) for c in data.get("key_concepts", [])][:10]
            entities = [str(e) for e in data.get("entities", [])][:10]
            steps.append("Sommario strutturato generato")
            return summary, key_concepts, entities

        except Exception as e:
            steps.append(f"Sommario auto-generato ({e})")
            words = raw_text.split()
            summary = " ".join(words[:80]) + "..." if len(words) > 80 else raw_text
            return summary, [], []

    # -------------------------------------------------------------------------
    # Context builder
    # -------------------------------------------------------------------------

    def _build_structured_extract(
        self,
        filename: str,
        summary: str,
        key_concepts: list,
        entities: list,
        visual_descriptions: list,
        raw_text: str,
    ) -> str:
        parts = [f"## DOCUMENTO: {filename}"]

        if summary:
            parts.append(f"\n### Sommario\n{summary}")

        if key_concepts:
            parts.append(f"\n### Concetti chiave\n" + " | ".join(key_concepts))

        if entities:
            parts.append(f"\n### Entità principali\n" + " | ".join(entities))

        if visual_descriptions:
            parts.append("\n### Contenuto visivo")
            for desc in visual_descriptions:
                parts.append(desc)

        if raw_text:
            preview = raw_text[:MAX_STRUCTURED_PREVIEW_CHARS]
            parts.append(f"\n### Contenuto testuale\n{preview}")
            if len(raw_text) > MAX_STRUCTURED_PREVIEW_CHARS:
                parts.append(f"\n[...testo troncato ({len(raw_text)} caratteri totali)]")

        return "\n".join(parts)


document_processor = DocumentProcessor()
